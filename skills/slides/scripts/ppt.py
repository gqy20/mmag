"""Deterministic PPTX renderer for the governed slides Skill."""

from __future__ import annotations

import argparse
import json
from pathlib import Path
from typing import Any

from pptx import Presentation
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt

MAX_SLIDES = 40
MAX_TEXT = 4000


def _text(value: Any, field: str, *, allow_empty: bool = False) -> str:
    if not isinstance(value, str) or (not allow_empty and not value.strip()):
        raise ValueError(f"{field} must be a non-empty string")
    if len(value) > MAX_TEXT:
        raise ValueError(f"{field} exceeds its length limit")
    return value.strip()


def _deck(payload: dict[str, Any]) -> dict[str, Any]:
    deck = payload.get("deck")
    if not isinstance(deck, dict):
        raise ValueError("input must contain a deck object")
    slides = deck.get("slides")
    if not isinstance(slides, list) or not 1 <= len(slides) <= MAX_SLIDES:
        raise ValueError("deck slides must contain between 1 and 40 items")
    return {
        "title": _text(deck.get("title"), "title"),
        "audience": _text(deck.get("audience"), "audience"),
        "objective": _text(deck.get("objective"), "objective"),
        "narrative": _text(deck.get("narrative"), "narrative"),
        "slides": [_slide(item, index) for index, item in enumerate(slides, 1)],
    }


def _slide(value: Any, index: int) -> dict[str, Any]:
    if not isinstance(value, dict):
        raise ValueError(f"slide {index} must be an object")
    body = value.get("body", ())
    if not isinstance(body, list) or len(body) > 12:
        raise ValueError(f"slide {index} body must be an array with at most 12 items")
    return {
        "title": _text(value.get("title"), f"slide {index} title"),
        "purpose": _text(value.get("purpose"), f"slide {index} purpose"),
        "body": [_text(item, f"slide {index} body") for item in body],
        "notes": _text(value.get("notes", ""), f"slide {index} notes", allow_empty=True),
    }


def render(input_path: Path, output_path: Path) -> None:
    payload = json.loads(input_path.read_text(encoding="utf-8"))
    if not isinstance(payload, dict):
        raise ValueError("input must be a JSON object")
    deck = _deck(payload)
    presentation = Presentation()
    presentation.slide_width = Inches(13.333)
    presentation.slide_height = Inches(7.5)
    _title_slide(presentation, deck)
    for slide in deck["slides"]:
        _content_slide(presentation, slide)
    output_path.parent.mkdir(parents=True, exist_ok=True)
    presentation.save(output_path)


def _title_slide(presentation: Presentation, deck: dict[str, Any]) -> None:
    slide = presentation.slides.add_slide(presentation.slide_layouts[0])
    slide.shapes.title.text = deck["title"]
    subtitle = slide.placeholders[1]
    subtitle.text = f"{deck['objective']}\nAudience: {deck['audience']}"
    slide.shapes.title.text_frame.paragraphs[0].font.size = Pt(30)


def _content_slide(presentation: Presentation, content: dict[str, Any]) -> None:
    slide = presentation.slides.add_slide(presentation.slide_layouts[1])
    slide.shapes.title.text = content["title"]
    title = slide.shapes.title.text_frame.paragraphs[0]
    title.font.size = Pt(26)
    title.font.bold = True
    body = slide.placeholders[1].text_frame
    body.clear()
    points = content["body"] or [content["purpose"]]
    for index, text in enumerate(points):
        paragraph = body.paragraphs[0] if index == 0 else body.add_paragraph()
        paragraph.text = text
        paragraph.font.size = Pt(18)
        paragraph.alignment = PP_ALIGN.LEFT
        paragraph.space_after = Pt(10)
    slide.notes_slide.notes_text_frame.text = content["notes"]


def main() -> None:
    parser = argparse.ArgumentParser()
    parser.add_argument("action", choices=("render",))
    parser.add_argument("--input", required=True, type=Path)
    parser.add_argument("--output", required=True, type=Path)
    arguments = parser.parse_args()
    render(arguments.input, arguments.output)


if __name__ == "__main__":
    main()
